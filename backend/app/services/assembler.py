import copy
import tempfile
from pathlib import Path
from pptx import Presentation


def assemble_deck_from_slides(slide_records, storage_dir: Path) -> Path:
    """
    Assemble a new .pptx from selected slides.

    slide_records: list of Slide ORM objects (with deck_id, slide_index, slide_data_path)
    storage_dir: base storage directory

    Returns: Path to the temporary .pptx file
    """
    if not slide_records:
        raise ValueError("No slides to assemble")

    # Group slides by source deck file
    deck_groups = {}
    for i, slide in enumerate(slide_records):
        source_path = str(storage_dir / slide.slide_data_path)
        if source_path not in deck_groups:
            deck_groups[source_path] = []
        deck_groups[source_path].append({
            "slide_index": slide.slide_index,
            "output_order": i,
        })

    # Create output presentation using first source's dimensions
    first_source_path = str(storage_dir / slide_records[0].slide_data_path)
    first_prs = Presentation(first_source_path)

    output = Presentation()
    output.slide_width = first_prs.slide_width
    output.slide_height = first_prs.slide_height

    # Collect slides in output order
    ordered_items = [None] * len(slide_records)

    for source_path, items in deck_groups.items():
        prs = Presentation(source_path)
        for item in items:
            idx = item["slide_index"]
            if idx < len(prs.slides):
                ordered_items[item["output_order"]] = (prs, idx)

    # Add each slide to output
    for item in ordered_items:
        if item is None:
            continue
        prs, slide_idx = item
        source_slide = prs.slides[slide_idx]
        _copy_slide(source_slide, output)

    # Save to temp file
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    output.save(tmp.name)
    tmp.close()

    return Path(tmp.name)


def _copy_slide(source_slide, target_prs):
    """
    Copy a single slide into the target presentation.
    Uses blank layout + deep copy of cSld (common slide data).
    """
    # Find a blank layout
    layout = None
    for sl in target_prs.slide_layouts:
        if sl.name.lower() in ("blank", "leer", "vide", "vuoto"):
            layout = sl
            break
    if layout is None:
        layout = target_prs.slide_layouts[len(target_prs.slide_layouts) - 1]

    new_slide = target_prs.slides.add_slide(layout)

    # Remove default placeholder shapes
    for ph in list(new_slide.placeholders):
        elem = ph._element
        elem.getparent().remove(elem)

    # Copy the cSld (common slide data) content
    source_cSld = source_slide._element.find(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld'
    )
    target_cSld = new_slide._element.find(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}cSld'
    )

    if source_cSld is not None and target_cSld is not None:
        new_cSld = copy.deepcopy(source_cSld)
        parent = target_cSld.getparent()
        parent.replace(target_cSld, new_cSld)

    # Copy image and other media relationships
    for rel_key, rel in source_slide.part.rels.items():
        if "image" in rel.reltype or "media" in rel.reltype:
            try:
                new_slide.part.relate_to(
                    rel.target_part, rel.reltype, rel_key, is_external=False
                )
            except Exception:
                pass

    return new_slide
